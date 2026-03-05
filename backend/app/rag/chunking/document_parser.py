"""Multi-format document parser (inspired by RAGFlow deep document understanding)."""

import io
import csv
import json
from pathlib import Path

import chardet
from bs4 import BeautifulSoup
from markdownify import markdownify


class DocumentParser:
    """Parse various document formats into plain text for chunking."""

    def parse(self, content: bytes, filename: str) -> str:
        ext = Path(filename).suffix.lower()
        parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_xlsx,
            ".csv": self._parse_csv,
            ".html": self._parse_html,
            ".htm": self._parse_html,
            ".json": self._parse_json,
            ".md": self._parse_text,
            ".txt": self._parse_text,
        }
        parser = parsers.get(ext, self._parse_text)
        return parser(content)

    def _parse_pdf(self, content: bytes) -> str:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"[Page {page_num}]\n{page_text}")
        return "\n\n".join(text_parts)

    def _parse_docx(self, content: bytes) -> str:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading structure
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        hashes = "#" * int(level)
                    except ValueError:
                        hashes = "#"
                    parts.append(f"{hashes} {para.text}")
                else:
                    parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts)

    def _parse_pptx(self, content: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts)

    def _parse_xlsx(self, content: bytes) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    def _parse_csv(self, content: bytes) -> str:
        encoding = chardet.detect(content)["encoding"] or "utf-8"
        text = content.decode(encoding)
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            rows.append(" | ".join(row))
        return "\n".join(rows)

    def _parse_html(self, content: bytes) -> str:
        encoding = chardet.detect(content)["encoding"] or "utf-8"
        html = content.decode(encoding)
        return markdownify(html, heading_style="ATX", strip=["img", "script", "style"])

    def _parse_json(self, content: bytes) -> str:
        encoding = chardet.detect(content)["encoding"] or "utf-8"
        data = json.loads(content.decode(encoding))
        return json.dumps(data, indent=2, ensure_ascii=False)

    def _parse_text(self, content: bytes) -> str:
        encoding = chardet.detect(content)["encoding"] or "utf-8"
        return content.decode(encoding)
